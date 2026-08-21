"""
Modul pemrosesan dokumen presentasi PowerPoint (.pptx / .ppt).

Mengekstrak slide presentasi menjadi teks Markdown terstruktur yang siap dichunking:
  - Menjaga judul slide (`## Slide N: [Judul]`)
  - Menjaga hierarki bullet points (poin-poin bertingkat)
  - Mengonversi tabel presentasi ke format Markdown Table (GFM)
  - Mengekstrak catatan pembicara (*speaker notes*)
  - Merender slide presentasi ke gambar resolusi tinggi untuk analisis VLM
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Batas lisensi Spire.Presentation Free: hanya 10 slide pertama per objek
# presentasi yang dirender penuh; slide ke-11 dst. menjadi blank + watermark.
# Karena itu rendering selalu dilakukan bertahap dalam batch sebesar nilai ini.
SPIRE_FREE_SLIDE_LIMIT: int = 10


def _table_to_markdown(table: Any) -> str:
    """Ubah objek tabel python-pptx menjadi teks tabel Markdown (GFM)."""
    rows: list[list[str]] = []
    for row in table.rows:
        cell_texts = [cell.text.replace("\n", " ").strip() for cell in row.cells]
        rows.append(cell_texts)

    if not rows:
        return ""

    header = rows[0]
    separator = ["---"] * len(header)
    md_lines: list[str] = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in rows[1:]:
        padded = row + [""] * (len(header) - len(row))
        md_lines.append("| " + " | ".join(padded[: len(header)]) + " |")

    return "\n".join(md_lines)


def count_presentation_slides(pptx_path: str | Path) -> int:
    """
    Hitung jumlah slide presentasi (.pptx / .ppt) menggunakan Spire.Presentation.
    """
    from spire.presentation import Presentation

    path_obj = Path(pptx_path).resolve()
    if not path_obj.exists():
        raise FileNotFoundError(f"File presentasi tidak ditemukan: {path_obj}")

    prs = Presentation()
    try:
        prs.LoadFromFile(str(path_obj))
        return prs.Slides.Count
    finally:
        prs.Dispose()


def _split_pptx_to_temp(
    src_path: Path,
    keep_start: int,
    keep_end: int,
    tmp_dir: Path,
) -> Path:
    """
    Buat salinan PPTX yang hanya mempertahankan slide indeks [keep_start, keep_end)
    (0-based) dan menyimpannya ke tmp_dir. Menggunakan python-pptx (tanpa batas lisensi).

    Slide di luar rentang dihapus melalui manipulasi `sldIdLst` + `drop_rel`,
    lalu disimpan sebagai file PPTX bersih (maksimal 10 slide) untuk di-render Spire.
    """
    from pptx import Presentation

    prs = Presentation(str(src_path))
    total = len(prs.slides)
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)

    # Hapus slide di luar rentang (urutkan descending agar indeks tetap valid).
    to_remove = [i for i in range(total) if not (keep_start <= i < keep_end)]
    for i in sorted(to_remove, reverse=True):
        prs.part.drop_rel(slides[i].rId)
        sldIdLst.remove(slides[i])

    out_path = tmp_dir / f"batch_{keep_start}_{keep_end}.pptx"
    prs.save(str(out_path))
    return out_path


def render_presentation_slides_to_images(
    pptx_path: str | Path,
    output_dir: str | Path | None = None,
    slides: list[int] | None = None,
    batch_size: int = SPIRE_FREE_SLIDE_LIMIT,
) -> list[Path]:
    """
    Render kanvas slide PowerPoint menjadi file gambar PNG (satu gambar per slide kanvas)
    dengan pemrosesan BERTAHAP per batch agar menghindari limitasi lisensi Spire.Presentation
    Free (maksimal 10 slide per objek presentasi; slide ke-11 dst. akan blank + watermark).

    Strategi: PPTX dipotong per batch (maksimal 10 slide) menggunakan python-pptx
    (tanpa batas lisensi) menjadi file PPTX kecil, lalu masing-masing file kecil
    di-render oleh Spire — sehingga Spire tidak pernah melihat presentasi > 10 slide.
    File .ppt (format lama) tidak bisa dipotong oleh python-pptx, sehingga menggunakan
    fallback AppendBySlide ke objek presentasi baru.

    Args:
        pptx_path: Path file presentasi (.pptx / .ppt).
        output_dir: Direktori penyimpanan gambar (default: <folder_pptx>/<stem>_slides).
        slides: Daftar indeks slide 0-based yang ingin dirender (None = seluruh slide).
        batch_size: Jumlah slide per batch render (default 10 = limit lisensi Free;
            jangan dinaikkan melebihi 10 pada versi Free).

    Nama file: `slide_<N>.png` (N mulai dari 1, sesuai nomor slide asli).
    Mengembalikan daftar path gambar yang dihasilkan (berurutan sesuai indeks input).
    """
    import tempfile

    from spire.presentation import Presentation

    path_obj = Path(pptx_path).resolve()
    if not path_obj.exists():
        raise FileNotFoundError(f"File presentasi tidak ditemukan: {path_obj}")

    out_dir = (
        Path(output_dir).resolve()
        if output_dir
        else (path_obj.parent / f"{path_obj.stem}_slides").resolve()
    )
    out_dir.mkdir(parents=True, exist_ok=True)

    is_pptx = path_obj.suffix.lower() == ".pptx"

    # Hitung total slide: python-pptx untuk .pptx, Spire untuk .ppt.
    if is_pptx:
        from pptx import Presentation as PptxPresentation

        probe = PptxPresentation(str(path_obj))
        total = len(probe.slides)
    else:
        src_probe = Presentation()
        src_probe.LoadFromFile(str(path_obj))
        total = src_probe.Slides.Count
        src_probe.Dispose()

    if total == 0:
        return []

    if slides is None:
        target_indices = list(range(total))
    else:
        target_indices = [i for i in slides if 0 <= i < total]
    if not target_indices:
        return []

    # Clamp: versi Free tidak bisa merender > 10 slide per objek presentasi.
    effective_batch = max(1, min(batch_size, SPIRE_FREE_SLIDE_LIMIT))

    generated_images: list[Path] = []
    with tempfile.TemporaryDirectory(prefix="pptx_batch_") as tmp_dir_name:
        tmp_dir = Path(tmp_dir_name)
        for b_start in range(0, len(target_indices), effective_batch):
            chunk = target_indices[b_start : b_start + effective_batch]

            batch_prs = Presentation()
            try:
                if is_pptx:
                    # Potong PPTX menjadi file kecil (hanya slide di chunk), lalu render.
                    split_path = _split_pptx_to_temp(
                        path_obj, chunk[0], chunk[-1] + 1, tmp_dir
                    )
                    batch_prs.LoadFromFile(str(split_path))
                else:
                    # Fallback .ppt: AppendBySlide ke presentasi baru.
                    # python-pptx tidak bisa membaca .ppt, jadi load source per batch.
                    src = Presentation()
                    src.LoadFromFile(str(path_obj))
                    try:
                        batch_prs.Slides.RemoveAt(0)
                        try:
                            batch_prs.SlideSize.Size = src.SlideSize.Size
                        except Exception:  # noqa: BLE001, S110 - fallback ke ukuran default
                            pass
                        for idx in chunk:
                            batch_prs.Slides.AppendBySlide(src.Slides[idx])
                    finally:
                        src.Dispose()

                for pos, src_idx in enumerate(chunk):
                    image = batch_prs.Slides[pos].SaveAsImage()
                    out_img = (out_dir / f"slide_{src_idx + 1}.png").resolve()
                    image.Save(str(out_img))
                    if out_img.exists():
                        generated_images.append(out_img)
            finally:
                batch_prs.Dispose()

    return generated_images


def pptx_to_structured_text(pptx_path: str | Path) -> list[dict[str, Any]]:
    """
    Ekstrak presentasi PPTX menjadi list struktur per slide.
    """
    path_obj = Path(pptx_path)
    if not path_obj.exists():
        raise FileNotFoundError(f"File presentasi tidak ditemukan: {path_obj}")

    prs = Presentation(str(path_obj))
    slides_data: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title: str = ""
        body_lines: list[str] = []
        notes_text: str = ""
        image_count: int = 0

        # 1. Ambil judul slide jika ada
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip()

        # 2. Iterasi shape dalam slide
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            # A. Gambar & Media
            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA):
                image_count += 1
                name = getattr(shape, "name", f"Image_{image_count}")
                body_lines.append(f"*[Visual / Diagram: {name}]*")

            # B. Tabel
            elif shape.has_table:
                md_table = _table_to_markdown(shape.table)
                if md_table:
                    body_lines.append(md_table)

            # C. Teks & Bullet points
            elif shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    text = p.text.strip()
                    if not text:
                        continue
                    if not slide_title and not body_lines:
                        slide_title = text
                        continue

                    level: int = getattr(p, "level", 0)
                    indent: str = "  " * level
                    body_lines.append(f"{indent}- {text}")

        # 3. Ambil speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            nt = slide.notes_slide.notes_text_frame.text.strip()
            if nt:
                notes_text = nt

        # 4. Susun markdown per slide
        title_display = slide_title or "Tanpa Judul"
        md_content_lines = [f"## Slide {idx}: {title_display}\n"]
        if body_lines:
            md_content_lines.append("\n".join(body_lines))
        if notes_text:
            md_content_lines.append(f"\n> **Speaker Notes:** {notes_text}")

        slide_markdown = "\n".join(md_content_lines)

        slides_data.append(
            {
                "slide_number": idx,
                "title": title_display,
                "markdown": slide_markdown,
                "notes": notes_text,
                "image_count": image_count,
            }
        )

    return slides_data


def process_presentation(pptx_path: str | Path) -> str:
    """
    Ekstrak seluruh file PPTX menjadi satu dokumen Markdown terpadu siap chunking.
    """
    slides = pptx_to_structured_text(pptx_path)
    file_stem = Path(pptx_path).stem.replace("_", " ").title()

    doc_lines: list[str] = [f"# {file_stem}\n"]
    for s in slides:
        doc_lines.append(s["markdown"])
        doc_lines.append("\n---\n")

    return "\n".join(doc_lines).strip()
